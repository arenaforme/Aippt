"""
PDF 转 PPTX 转换器
将 PDF 演示文稿转换为可编辑的 PPTX 文件

支持两种模式：
1. 矢量 PDF：直接提取文本和图片
2. 图片 PDF：使用 OCR 识别文字（如 NotebookLM 导出的 PDF）
"""

import io
import logging
import tempfile
from pathlib import Path
from typing import Optional, Callable, Union

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .models import PDFConversionResult, PDFPageData, PDFTextBlock, PDFImage
from .parser import PDFParser

logger = logging.getLogger(__name__)


class PDFConverter:
    """PDF 转 PPTX 转换器"""

    # 标准 16:9 幻灯片尺寸（EMU）
    SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
    SLIDE_HEIGHT = Emu(6858000)  # 7.5 inches

    def __init__(
        self,
        baidu_api_key: str = None,
        baidu_secret_key: str = None
    ):
        """
        初始化转换器

        Args:
            baidu_api_key: 百度 OCR API Key（用于图片 PDF）
            baidu_secret_key: 百度 OCR Secret Key
        """
        self.baidu_api_key = baidu_api_key
        self.baidu_secret_key = baidu_secret_key
        self.prs: Optional[Presentation] = None

    def convert(
        self,
        pdf_path: Union[str, Path],
        output_path: Union[str, Path],
        progress_callback: Optional[Callable[[dict], None]] = None
    ) -> PDFConversionResult:
        """
        转换 PDF 为 PPTX
        自动检测 PDF 类型（矢量/图片），选择合适的转换方式
        """
        pdf_path = Path(pdf_path)
        output_path = Path(output_path)

        if not pdf_path.exists():
            return PDFConversionResult(
                success=False,
                error_message=f"PDF 文件不存在: {pdf_path}"
            )

        try:
            # 检测 PDF 类型
            is_image_pdf = self._is_image_based_pdf(pdf_path)

            if is_image_pdf:
                logger.info(f"检测到图片型 PDF，使用 OCR 模式转换")
                return self._convert_image_pdf(pdf_path, output_path, progress_callback)
            else:
                logger.info(f"检测到矢量型 PDF，使用直接提取模式")
                return self._convert_vector_pdf(pdf_path, output_path, progress_callback)

        except Exception as e:
            logger.exception(f"PDF 转换失败: {e}")
            return PDFConversionResult(
                success=False,
                error_message=str(e)
            )

    def _is_image_based_pdf(self, pdf_path: Path) -> bool:
        """
        检测 PDF 是否为图片型（每页是一张大图片，无矢量文本）
        """
        with PDFParser(pdf_path) as parser:
            page_count = parser.get_page_count()
            if page_count == 0:
                return False

            # 检查前几页
            check_pages = min(3, page_count)
            image_pages = 0

            for i in range(check_pages):
                page_data = parser.parse_page(i)

                # 如果没有文本块，但有图片，认为是图片页
                if len(page_data.text_blocks) == 0 and len(page_data.images) > 0:
                    # 检查图片是否覆盖大部分页面
                    for img in page_data.images:
                        x0, y0, x1, y1 = img.bbox
                        img_area = (x1 - x0) * (y1 - y0)
                        page_area = page_data.width * page_data.height
                        if img_area > page_area * 0.5:  # 图片占页面 50% 以上
                            image_pages += 1
                            break

            # 如果大部分检查页都是图片页，认为是图片型 PDF
            return image_pages >= check_pages * 0.5

    def _convert_image_pdf(
        self,
        pdf_path: Path,
        output_path: Path,
        progress_callback: Optional[Callable[[dict], None]] = None
    ) -> PDFConversionResult:
        """
        转换图片型 PDF（使用 OCR）
        """
        # 检查 OCR 配置
        if not self.baidu_api_key or not self.baidu_secret_key:
            return PDFConversionResult(
                success=False,
                error_message="图片型 PDF 需要配置百度 OCR API（请在设置中配置 BAIDU_OCR_API_KEY 和 BAIDU_OCR_SECRET_KEY）"
            )

        # 导入 PPTConverter（延迟导入避免循环依赖）
        from services.ppt_converter import PPTConverter
        from services.ppt_converter.text_corrector import load_reference_text

        try:
            # 1. 从 PDF 提取图片到临时目录
            temp_dir = tempfile.mkdtemp(prefix="pdf_ocr_")
            image_paths = []

            with PDFParser(pdf_path) as parser:
                page_count = parser.get_page_count()

                if progress_callback:
                    progress_callback({
                        'current_page': 0,
                        'total': page_count,
                        'stage': 'extracting',
                        'stage_name': '正在从 PDF 提取图片...'
                    })

                for page_num in range(page_count):
                    # 将 PDF 页面渲染为图片（使用较低 DPI 避免超出 OCR 尺寸限制）
                    image_path = Path(temp_dir) / f"page_{page_num + 1:03d}.png"
                    parser.render_page_to_image(page_num, image_path, dpi=100)
                    image_paths.append(image_path)

                    if progress_callback:
                        progress_callback({
                            'current_page': page_num + 1,
                            'total': page_count,
                            'stage': 'extracting',
                            'stage_name': f'提取第 {page_num + 1} 页图片...'
                        })

            # 2. 加载参考文本（如果存在同名 .txt 文件）
            reference_text = load_reference_text(pdf_path)
            if reference_text:
                logger.info(f"已加载参考文本，将用于校正 OCR 结果")

            # 3. 使用 PPTConverter 进行 OCR 转换
            ppt_converter = PPTConverter(
                api_key=self.baidu_api_key,
                secret_key=self.baidu_secret_key,
                reference_text=reference_text
            )

            result = ppt_converter.convert_images(
                image_paths=image_paths,
                output_path=output_path,
                remove_text=True,  # 移除原图中的文字区域
                progress_callback=progress_callback
            )

            # 4. 清理临时文件
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)

            return PDFConversionResult(
                success=result.success,
                output_path=output_path if result.success else None,
                pages_count=result.slides_count,
                text_blocks_count=result.text_blocks_count,
                images_count=0,
                error_message=None if result.success else "OCR 转换失败"
            )

        except Exception as e:
            logger.exception(f"图片型 PDF 转换失败: {e}")
            return PDFConversionResult(
                success=False,
                error_message=str(e)
            )

    def _convert_vector_pdf(
        self,
        pdf_path: Path,
        output_path: Path,
        progress_callback: Optional[Callable[[dict], None]] = None
    ) -> PDFConversionResult:
        """
        转换矢量型 PDF（直接提取文本）
        """
        try:
            # 初始化 Presentation
            self.prs = Presentation()
            self.prs.slide_width = self.SLIDE_WIDTH
            self.prs.slide_height = self.SLIDE_HEIGHT

            total_text_blocks = 0
            total_images = 0

            with PDFParser(pdf_path) as parser:
                page_count = parser.get_page_count()

                if page_count == 0:
                    return PDFConversionResult(
                        success=False,
                        error_message="PDF 文件为空"
                    )

                for page_num in range(page_count):
                    if progress_callback:
                        progress_callback({
                            'current_page': page_num + 1,
                            'total': page_count,
                            'stage': 'parsing',
                            'stage_name': f'解析第 {page_num + 1} 页...'
                        })

                    page_data = parser.parse_page(page_num)
                    self._add_slide(page_data)

                    total_text_blocks += len(page_data.text_blocks)
                    total_images += len(page_data.images)

                    if progress_callback:
                        progress_callback({
                            'current_page': page_num + 1,
                            'total': page_count,
                            'completed': page_num + 1,
                            'stage': 'page_done',
                            'stage_name': f'第 {page_num + 1} 页完成',
                            'text_blocks_count': len(page_data.text_blocks)
                        })

            if progress_callback:
                progress_callback({
                    'current_page': page_count,
                    'total': page_count,
                    'completed': page_count,
                    'stage': 'generating',
                    'stage_name': '正在生成 PPTX 文件...'
                })

            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(str(output_path))

            return PDFConversionResult(
                success=True,
                output_path=output_path,
                pages_count=page_count,
                text_blocks_count=total_text_blocks,
                images_count=total_images
            )

        except Exception as e:
            logger.exception(f"矢量型 PDF 转换失败: {e}")
            return PDFConversionResult(
                success=False,
                error_message=str(e)
            )

    def _add_slide(self, page_data: PDFPageData) -> None:
        """添加一页幻灯片"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        scale_x = self.SLIDE_WIDTH / page_data.width
        scale_y = self.SLIDE_HEIGHT / page_data.height

        for img in page_data.images:
            self._add_image(slide, img, scale_x, scale_y)

        for text_block in page_data.text_blocks:
            self._add_text_box(slide, text_block, scale_x, scale_y)

    def _add_image(self, slide, img: PDFImage, scale_x: float, scale_y: float) -> None:
        """添加图片到幻灯片"""
        x0, y0, x1, y1 = img.bbox
        left = int(x0 * scale_x)
        top = int(y0 * scale_y)
        width = int((x1 - x0) * scale_x)
        height = int((y1 - y0) * scale_y)

        try:
            image_stream = io.BytesIO(img.image_data)
            slide.shapes.add_picture(image_stream, left, top, width, height)
        except Exception as e:
            logger.warning(f"添加图片失败: {e}")

    def _add_text_box(self, slide, text_block: PDFTextBlock, scale_x: float, scale_y: float) -> None:
        """添加文本框到幻灯片"""
        x0, y0, x1, y1 = text_block.bbox
        left = int(x0 * scale_x)
        top = int(y0 * scale_y)
        width = int((x1 - x0) * scale_x)
        height = int((y1 - y0) * scale_y * 1.2)

        font_size = int(text_block.font_size * (scale_y / Emu(914400)))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False

        p = tf.paragraphs[0]
        p.text = text_block.text
        p.alignment = PP_ALIGN.LEFT

        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(max(8, min(font_size, 72)))

        r, g, b = text_block.color
        run.font.color.rgb = RGBColor(r, g, b)

        if text_block.is_bold:
            run.font.bold = True
