"""
PPT 生成模块
根据识别结果生成可编辑的 .pptx 文件
"""

import logging
from pathlib import Path
from typing import Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from .models import SlideData, TextBlock

logger = logging.getLogger(__name__)


def _estimate_text_width(text: str, font_size_pt: float) -> float:
    """
    估算文字渲染后的宽度（EMU单位）

    中文字符宽度约等于字号，英文/数字约为字号的0.5-0.6倍
    """
    # 1 pt = 12700 EMU
    pt_to_emu = 12700

    width = 0
    for char in text:
        if '\u4e00' <= char <= '\u9fff':
            # 中文字符：宽度约等于字号
            width += font_size_pt * pt_to_emu
        elif char.isalpha():
            # 英文字母：宽度约为字号的0.55倍
            width += font_size_pt * pt_to_emu * 0.55
        elif char.isdigit():
            # 数字：宽度约为字号的0.55倍
            width += font_size_pt * pt_to_emu * 0.55
        else:
            # 其他字符（标点等）：宽度约为字号的0.4倍
            width += font_size_pt * pt_to_emu * 0.4

    return width


def _calculate_optimal_font_size(
    text: str,
    bbox_width_emu: int,
    original_font_size: int
) -> int:
    """
    计算适合 bbox 宽度的最优字号

    如果原始字号会导致文字溢出，则缩小字号
    """
    # 先用原始字号估算文字宽度
    estimated_width = _estimate_text_width(text, original_font_size)

    text_preview = text[:20] if len(text) > 20 else text
    logger.info(
        f"字号计算: text='{text_preview}', "
        f"bbox_width={bbox_width_emu}, estimated_width={estimated_width:.0f}, "
        f"original_size={original_font_size}"
    )

    # 如果估算宽度小于 bbox 宽度，使用原始字号
    if estimated_width <= bbox_width_emu:
        logger.info(f"  -> 无需缩小，使用原始字号 {original_font_size}pt")
        return original_font_size

    # 否则，按比例缩小字号
    ratio = bbox_width_emu / estimated_width
    adjusted_size = int(original_font_size * ratio * 0.95)  # 留5%余量

    # 字号下限为8pt
    final_size = max(8, adjusted_size)
    logger.info(
        f"  -> 字号缩小: {original_font_size}pt -> {final_size}pt "
        f"(ratio={ratio:.2f})"
    )
    return final_size


class PPTGenerator:
    """PPT 生成器"""

    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.SLIDE_WIDTH_INCHES)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT_INCHES)

    def add_slide(self, slide_data: SlideData) -> None:
        """添加一页幻灯片"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        self._add_background_image(slide, slide_data)

        for text_block in slide_data.text_blocks:
            self._add_text_box(slide, text_block, slide_data)

    def _add_background_image(self, slide, slide_data: SlideData) -> None:
        """添加背景图片"""
        if not slide_data.image_path.exists():
            return

        slide.shapes.add_picture(
            str(slide_data.image_path),
            Inches(0),
            Inches(0),
            width=self.prs.slide_width,
            height=self.prs.slide_height
        )

    def _add_text_box(
        self,
        slide,
        text_block: TextBlock,
        slide_data: SlideData
    ) -> None:
        """添加文本框"""
        x, y, w, h = text_block.bbox

        scale_x = self.prs.slide_width / slide_data.width
        scale_y = self.prs.slide_height / slide_data.height

        left = int(x * scale_x)
        top = int(y * scale_y)
        width = int(w * scale_x)
        height = int(h * scale_y * 1.3)  # 高度留一些余量

        # 确保文本框不超出幻灯片右边界
        max_width = self.prs.slide_width - left
        if width > max_width:
            width = max_width

        # 计算最优字号：如果文字超出 bbox 宽度，则缩小字号
        font_size = _calculate_optimal_font_size(
            text_block.text,
            width,
            text_block.font_size
        )

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False  # 禁用自动换行，保持单行

        p = tf.paragraphs[0]
        p.text = text_block.text
        p.alignment = PP_ALIGN.LEFT

        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_size)
        run.font.name = text_block.font_name

        # 设置东亚字体（中文字符使用此字体）
        run.font._element.set(qn('w:eastAsia'), text_block.font_name)
        # 同时设置 latin 和 ea 属性确保中文正确显示
        rPr = run.font._element
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            from lxml import etree
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', text_block.font_name)

        r, g, b = text_block.color
        run.font.color.rgb = RGBColor(r, g, b)

    def save(self, output_path: Union[str, Path]) -> Path:
        """保存 PPT 文件"""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        return output_path
